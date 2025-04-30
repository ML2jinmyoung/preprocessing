import os
import PyPDF2
import pptx
import openpyxl
import subprocess
import tempfile
import logging
import gc
from pathlib import Path
from typing import Optional, Dict, Callable, Union, List
from dataclasses import dataclass
from pdf2image import convert_from_path
from PIL import Image
import fitz  # PyMuPDF
import traceback

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

DEFAULT_CONFIG = {
    'input_dir': "./demo",
    'output_dir': "./demo_result",
    'min_text_length': 100,
    'pdf_dpi': 300,
    'pdf_zoom': 4,
    'image_min_dpi': 200,
    'image_target_dpi': 300,
    'chunk_size': 1000,
    'overlap_size': 50,
    'korean_threshold': 0.3,  # 한국어 비율 임계값 (더 낮게 설정)
    'ocr_reset_interval': 10  # 이 횟수만큼 OCR을 수행하면 OCR 엔진을 재생성
}

SUPPORTED_FORMATS = {
    'image': ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif'],
    'document': ['.pdf', '.doc', '.docx'],
    # 'presentation': ['.ppt', '.pptx'],
    # 'spreadsheet': ['.xls', '.xlsx']
}

DOCUMENT_GROUP = 'LAW'

@dataclass
class ConversionResult:
    success: bool
    content: str
    error: Optional[str] = None
    is_korean: bool = True  # 한국어 문서 여부를 저장하는 필드 추가

class DocumentConverter:
    def __init__(self, config: dict = None):
        self.config = config or DEFAULT_CONFIG
        self._setup_converters()
        self.failed_files = []
        self.non_korean_files = []  # 한국어가 아닌 문서 목록 추가
        self.ocr_count = 0  # OCR 수행 횟수를 추적하는 카운터
        self.ocr_instance = None  # OCR 인스턴스를 필요할 때만 생성
        
    def _get_ocr_instance(self):
        """
        OCR 인스턴스를 얻거나 필요시 재생성
        """
        if self.ocr_instance is None:
            from ocr import get_ocr_instance, reset_ocr_instance
            self.ocr_instance = get_ocr_instance()
            self.ocr_count = 0
            logger.info("새 OCR 인스턴스 생성됨")
        
        self.ocr_count += 1
        
        # 일정 횟수 이상 사용했다면 인스턴스 재생성 준비
        if self.ocr_count >= self.config['ocr_reset_interval']:
            reset_ocr_instance()
            self.ocr_instance = None
            gc.collect()  # 가비지 컬렉션 강제 실행
            
        return self.ocr_instance

    def _setup_converters(self) -> None:
        """
        파일 형식별 변환 함수 매핑
        """
        self.converters = {
            **{ext: self._convert_image_to_txt for ext in SUPPORTED_FORMATS['image']},
            '.pdf': self._convert_pdf_to_txt,
            **{ext: self._convert_doc_to_txt for ext in ['.doc', '.docx']},
            **{ext: self._convert_pptx_to_txt for ext in ['.ppt', '.pptx']},
            **{ext: self._convert_xlsx_to_txt for ext in ['.xls', '.xlsx']}
        }

    def _check_if_korean(self, text: str) -> bool:
        """
        is_korean_dominant 함수를 사용하여 한국어 문서인지 확인
        """
        if not text or len(text.strip()) == 0:
            return False
            
        from ocr import is_korean_dominant
        return is_korean_dominant(text, self.config['korean_threshold'])

    def _scale_image_if_needed(self, image: Image.Image) -> Image.Image:
        """
        이미지 DPI 조정
        """
        try:
            dpi = image.info.get('dpi', (72, 72))
            if isinstance(dpi, tuple) and min(dpi) < self.config['image_min_dpi']:
                scale_factor = self.config['image_target_dpi'] / min(72, min(dpi))
                new_size = (int(image.width * scale_factor), int(image.height * scale_factor))
                return image.resize(new_size, Image.LANCZOS)
        except Exception as e:
            logger.warning(f"DPI 조정 실패: {e}")
        return image

    def _convert_image_to_txt(self, file_path: str) -> ConversionResult:
        """
        이미지 파일에서 텍스트 추출
        """
        try:
            image = Image.open(file_path).convert("RGB")
            image = self._scale_image_if_needed(image)
            
            # OCR 처리를 위한 함수 호출 (ocr.py 분리)
            from ocr import try_multiple_ocr_approaches
            text = try_multiple_ocr_approaches(image)
            
            # OCR 처리 후 메모리 정리
            image = None
            gc.collect()
            
            if not text.strip():
                return ConversionResult(False, "", "‼️ 텍스트를 추출할 수 없습니다")
                
            return ConversionResult(True, text)
        except Exception as e:
            logger.error(f"이미지 처리 중 에러: {str(e)}")
            traceback.print_exc()
            return ConversionResult(False, "", f"‼️ 이미지 처리 중 에러: {str(e)}")

    def _convert_pdf_to_txt(self, file_path: str) -> ConversionResult:
        """
        PDF 파일에서 텍스트 추출
        """
        try:
            text = self._extract_pdf_text(file_path)
            if len(text.strip()) < self.config['min_text_length']:
                text = self._process_scanned_pdf(file_path)
                
            if not text.strip():
                return ConversionResult(False, "", "‼️ 텍스트 추출 실패")
                
            # PDF OCR 결과에 대해서만 한국어 체크
            is_korean = self._check_if_korean(text)
            if not is_korean:
                logger.info(f"한국어가 아닌 PDF 문서 감지됨: {file_path}")
                return ConversionResult(True, text, is_korean=False)
                
            return ConversionResult(True, text)
        except Exception as e:
            logger.error(f"PDF 처리 중 에러: {str(e)}")
            traceback.print_exc()
            return ConversionResult(False, "", f"‼️ PDF 처리 중 에러: {str(e)}")

    def _extract_pdf_text(self, file_path: str) -> str:
        """
        일반 PDF에서 텍스트 추출
        """
        try:
            with open(file_path, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)
                return "\n".join(page.extract_text() for page in pdf_reader.pages)
        except Exception as e:
            logger.warning(f"PDF 텍스트 추출 실패: {e}")
            return ""

    def _process_scanned_pdf(self, file_path: str) -> str:
        """
        스캔된 PDF 처리
        """
        try:
            return self._process_with_pdf2image(file_path)
        except Exception as e:
            logger.warning(f"‼️ PDF2Image 실패, PyMuPDF 시도: {e}")
            return self._process_with_pymupdf(file_path)

    def _process_with_pdf2image(self, file_path: str) -> str:
        """
        PDF2Image를 사용한 PDF 처리
        """
        images = convert_from_path(
            file_path,
            dpi=self.config['pdf_dpi'],
            thread_count=4,
            use_cropbox=True,
            strict=False,
            grayscale=False
        )
        return self._process_pdf_images(images)

    def _process_with_pymupdf(self, file_path: str) -> str:
        """
        PyMuPDF를 사용한 PDF 처리
        """
        pdf_document = fitz.open(file_path)
        mat = fitz.Matrix(self.config['pdf_zoom'], self.config['pdf_zoom'])
        images = []
        
        for page in pdf_document:
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
            
        return self._process_pdf_images(images)

    def _process_pdf_images(self, images: list) -> str:
        """
        PDF 이미지들에서 텍스트 추출
        """
        text = ""
        for i, img in enumerate(images, 1):
            logger.info(f"페이지 {i}/{len(images)} OCR 처리 중...")
            
            # OCR 처리를 위한 함수 호출
            from ocr import try_multiple_ocr_approaches
            page_text = try_multiple_ocr_approaches(img)
            text += f"{page_text}\n\n"
            
            # 각 페이지 처리 후 메모리 정리
            img = None
            gc.collect()
            
        return text

    def _convert_doc_to_txt(self, file_path: str) -> ConversionResult:
        """
        DOC/DOCX 파일에서 텍스트 추출
        """
        ext = Path(file_path).suffix.lower()
        try:
            result = None
            if ext == '.docx':
                result = self._convert_docx(file_path)
            else:
                result = self._convert_doc(file_path)
                    
            return result
        except Exception as e:
            import traceback
            trace = traceback.format_exc()
            return ConversionResult(False, "", f"문서 처리 중 에러: {str(e)}{{trace}}")

    def _convert_docx(self, file_path: str) -> ConversionResult:
        """
        DOCX 파일 처리 (docx2txt 라이브러리 사용)
        """
        try:
            import docx2txt
            text = docx2txt.process(file_path)
            return ConversionResult(True, text)
        except Exception as e:
            import traceback
            tb = traceback.format_exc()
            return ConversionResult(False, "", f"DOCX 변환 중 오류 발생: {str(e)}\n{tb}")

    def _convert_doc(self, file_path: str) -> ConversionResult:
        """
        DOC 파일 처리 (실제 RTF 포맷인지 확인)
        """
        for tool in ['antiword', 'catdoc']:
            try:
                result = subprocess.run([tool, file_path], capture_output=True, text=True, check=True)
                return ConversionResult(True, result.stdout)
            except:
                continue
        return ConversionResult(False, "", "‼️ DOC 파일 처리 실패")

    def _convert_pptx_to_txt(self, file_path: str) -> ConversionResult:
        """
        PPT/PPTX 파일에서 텍스트 추출
        """
        try:
            prs = pptx.Presentation(file_path)
            text = "\n".join(
                shape.text for slide in prs.slides 
                for shape in slide.shapes if hasattr(shape, "text")
            )
            
            return ConversionResult(True, text)
        except Exception as e:
            try:
                result = subprocess.run(['pptx2txt', file_path], capture_output=True, text=True, check=True)
                return ConversionResult(True, result.stdout)
            except:
                return ConversionResult(False, "", f"‼️ 프레젠테이션 처리 중 에러: {str(e)}")

    def _convert_xlsx_to_txt(self, file_path: str) -> ConversionResult:
        """
        XLSX/XLS 파일에서 텍스트 추출
        """
        try:
            result = self._convert_xlsx_with_openpyxl(file_path)
            return result
        except Exception as e:
            logger.warning(f"‼️ OpenPyXL 실패, xlsx2csv 시도: {e}")
            result = self._convert_xlsx_with_xlsx2csv(file_path)
            return result

    def _convert_xlsx_with_openpyxl(self, file_path: str) -> ConversionResult:
        """
        OpenPyXL을 사용한 엑셀 처리
        """
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join(str(cell) if cell is not None else "" for cell in row) + "\n"
            text += "\n"
        return ConversionResult(True, text)

    def _convert_xlsx_with_xlsx2csv(self, file_path: str) -> ConversionResult:
        """
        xlsx2csv를 사용한 엑셀 처리
        """
        with tempfile.NamedTemporaryFile(suffix=".csv", delete=False) as tmp:
            try:
                subprocess.run(['xlsx2csv', file_path, tmp.name], check=True)
                with open(tmp.name, 'r', encoding='utf-8', errors='replace') as f:
                    return ConversionResult(True, f.read())
            finally:
                os.unlink(tmp.name)

    def convert_file(self, file_path: str) -> ConversionResult:
        """
        파일을 텍스트로 변환
        """
        ext = Path(file_path).suffix.lower()
        converter = self.converters.get(ext)
        
        if not converter:
            return self._handle_unknown_format(file_path)
            
        result = converter(file_path)
        
        # 파일 변환 후 메모리 정리
        gc.collect()
        
        return result

    def _handle_unknown_format(self, file_path: str) -> ConversionResult:
        """
        알 수 없는 형식의 파일 처리
        """
        try:
            result = subprocess.run(['file', '-b', file_path], capture_output=True, text=True, check=True)
            file_type = result.stdout.lower()
            
            if 'text' in file_type:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read()
                    return ConversionResult(True, text)
            
            return ConversionResult(False, "", f"‼️ 지원되지 않는 파일 형식: {Path(file_path).suffix}")
        except Exception as e:
            return ConversionResult(False, "", f"‼️ 파일 처리 중 에러: {str(e)}")

    def _display_progress(self, current: int, total: int, file_name: str) -> None:
        """
        진행 상황을 표시
        """
        progress = (current / total) * 100
        print(f"\r[{current}/{total}] ({progress:.1f}%) Processing: {file_name}", end="")

    def process_directory(self, input_dir: str, output_dir: str) -> None:
        """
        디렉토리 내의 모든 파일 처리
        """
        input_path = Path(input_dir)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        supported_exts = set()
        for ext_list in SUPPORTED_FORMATS.values():
            supported_exts.update(ext_list)
            
        # 처리할 모든 파일 목록 미리 수집
        all_files = [
            file_path for file_path in input_path.rglob('*')
            if file_path.is_file() and not file_path.name.startswith('.')
            and file_path.suffix.lower() in supported_exts
        ]
        
        total_files = len(all_files)
        print(f"\n📄 총 {total_files}개의 파일을 처리합니다.\n")
        
        try:
            # 배치 단위로 파일 처리
            batch_size = 5  # 한 번에 처리할 파일 수
            for i in range(0, total_files, batch_size):
                batch_files = all_files[i:i+batch_size]
                
                for idx, file_path in enumerate(batch_files):
                    current_file_num = i + idx + 1
                    try:
                        # _display_progress(current_file_num, total_files, file_path.name)
                        self._process_single_file(file_path, input_path, output_path)
                    except Exception as e:
                        logger.error(f"\n❌ 파일 처리 실패 ({file_path.name}): {str(e)}")
                        self.failed_files.append((str(file_path), str(e)))
                        continue
                    
                # 배치 처리 후 메모리 정리
                from ocr import reset_ocr_instance
                reset_ocr_instance()
                self.ocr_instance = None
                gc.collect()
            
                print()  # 새 줄로 이동
            
        except Exception as e:
            logger.error(f"\n❌ 디렉토리 처리 중 오류 발생: {str(e)}")
            raise

    def _process_single_file(self, file_path: Path, input_path: Path, output_path: Path) -> None:
        """
        단일 파일 처리
        """
        relative_path = file_path.relative_to(input_path.parent)
        logger.info(f"처리 중: {relative_path}")

        if file_path.suffix.lower() in SUPPORTED_FORMATS['image'] + ['.pdf']:
            logger.info("OCR 처리가 필요할 수 있어 시간이 오래 걸릴 수 있습니다...")

        result = self.convert_file(str(file_path))
        if result.success:
            # 한국어 문서인 경우에만 결과 저장
            if result.is_korean:
                self._save_result(result.content, file_path, relative_path, output_path)
            else:
                # 한국어가 아닌 경우 변환 실패 목록에 추가
                error_msg = "한국어가 아닌 문서"
                logger.warning(f"한국어가 아닌 문서 감지됨: {relative_path}")
                self.non_korean_files.append((str(relative_path), error_msg))
                self._save_result(result.content, file_path, relative_path, output_path)
        else:
            logger.error(f"변환 실패 ({relative_path}): {result.error}")
            self.failed_files.append((str(relative_path), result.error))

    def _chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> str:
        """
        텍스트를 chunk_size 만큼 나누고, overlap만큼 간격 주고 <Chunk> 태그로 감싸고 <Content>로 전체 감쌈
        """
        if chunk_size is None:
            chunk_size = self.config['chunk_size']
        if overlap is None:
            overlap = self.config['overlap_size']
            
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start += chunk_size - overlap  # 겹치도록 다음 시작 위치 설정

        chunked_text = "<Content>\n"
        for chunk in chunks:
            chunked_text += f"<Chunk>\n{chunk.strip()}\n</Chunk>\n\n"
        chunked_text += "</Content>\n"
        return chunked_text

    def _save_result(self, content: str, file_path: Path, relative_path: Path, output_path: Path) -> None:
        """
        변환 결과 저장
        """
        try:
            metadata = (
                "<Metadata>\n"
                f"    <Group>{DOCUMENT_GROUP}</Group>\n"
                f"    <File>{file_path.name}</File>\n"
                f"    <Path>{relative_path}</Path>\n"
                "</Metadata>\n\n"
            )
            chunked_body = self._chunk_text(content)

            document_text = "<Document>\n" + metadata + chunked_body + "</Document>\n"

            out_file = output_path / relative_path.parent / f"{file_path.stem}.txt"
            out_file.parent.mkdir(parents=True, exist_ok=True)

            with open(out_file, "w", encoding="utf-8") as f:
                f.write(document_text)

            chunk_count = chunked_body.count("<Chunk>")
            print(f"\n✅ 완료: {file_path.name} ({chunk_count} chunks)")
        except Exception as e:
            logger.error(f"\n❌ 저장 실패 ({out_file}): {e}")

def main():
    """
    메인 함수
    """
    converter = DocumentConverter()
    converter.process_directory(DEFAULT_CONFIG['input_dir'], DEFAULT_CONFIG['output_dir'])

    print_failure_list = False
    
    if converter.non_korean_files:
        print("\n=== ❌ 한국어가 아닌 파일 목록 ===")
        for path, error in converter.non_korean_files:
            print(f"- {path}: {error}")
        print_failure_list = True
    
    if converter.failed_files:
        print("\n=== ❌ 변환 실패 파일 목록 ===")
        for path, error in converter.failed_files:
            print(f"- {path}: {error}")
        print_failure_list = True
    
    if not print_failure_list:
        print("\n🎉 모든 파일이 성공적으로 처리되었습니다.")


if __name__ == "__main__":
    main()